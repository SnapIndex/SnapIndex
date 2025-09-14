import os
import mimetypes
from pypdf import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation
import csv
from image_classifier import ImageClassifier
import threading

# Global image classifier instance (lazy loaded)
_image_classifier = None
_image_classifier_failed = False
_image_classifier_lock = threading.Lock()

def get_image_classifier():
    """Get or create the global image classifier instance"""
    global _image_classifier, _image_classifier_failed
    if _image_classifier is not None:
        return _image_classifier
    if _image_classifier_failed:
        return None
    with _image_classifier_lock:
        if _image_classifier is not None:
            return _image_classifier
        if _image_classifier_failed:
            return None
        try:
            _image_classifier = ImageClassifier()
            print("✅ Image classifier initialized successfully")
        except Exception as e:
            print(f"⚠️ Failed to initialize image classifier: {e}")
            print("Images will be processed without AI classification")
            _image_classifier_failed = True
            _image_classifier = None
        return _image_classifier

# Supported file extensions
SUPPORTED_EXTENSIONS = {
    # PDF files
    '.pdf': 'pdf',
    
    # Microsoft Word documents
    '.docx': 'docx',
    '.doc': 'doc',
    
    # Microsoft Excel files
    '.xlsx': 'xlsx',
    '.xls': 'xls',
    '.csv': 'csv',
    
    # Microsoft PowerPoint presentations
    '.pptx': 'pptx',
    '.ppt': 'ppt',
    
    # Text files
    '.txt': 'txt',
    '.rtf': 'rtf',
    '.md': 'markdown',
    
    # Other text formats
    '.html': 'html',
    '.htm': 'html',
    '.xml': 'xml',
    '.json': 'json',
    '.log': 'log',
    '.ini': 'ini',
    '.cfg': 'cfg',
    '.conf': 'conf',
    
    # Image files
    '.jpg': 'image',
    '.jpeg': 'image',
    '.png': 'image',
    '.gif': 'image',
    '.bmp': 'image',
    '.tiff': 'image',
    '.tif': 'image',
    '.svg': 'image',
    '.webp': 'image',
    '.ico': 'image',
    '.psd': 'image',
    '.ai': 'image',
    '.eps': 'image',
    '.raw': 'image',
    '.cr2': 'image',
    '.nef': 'image',
    '.orf': 'image',
    '.sr2': 'image',
    '.arw': 'image',
    '.dng': 'image',
    '.heic': 'image',
    '.heif': 'image',
    '.avif': 'image',
    '.jfif': 'image',
    '.pbm': 'image',
    '.pgm': 'image',
    '.ppm': 'image',
    '.pnm': 'image',
    '.xbm': 'image',
    '.xpm': 'image',
    '.pcx': 'image',
    '.tga': 'image',
    '.exr': 'image',
    '.hdr': 'image',
    '.pic': 'image',
    '.pict': 'image'
}

def find_documents(folder_path):
    """Find all supported document files in the folder"""
    documents = []
    for root, _, files in os.walk(folder_path):
        for file in files:
            file_ext = os.path.splitext(file)[1].lower()
            if file_ext in SUPPORTED_EXTENSIONS:
                documents.append({
                    'path': os.path.join(root, file),
                    'name': file,
                    'extension': file_ext,
                    'type': SUPPORTED_EXTENSIONS[file_ext]
                })
    return documents

def extract_text_from_file(file_info):
    """Extract text from a single file based on its type"""
    file_path = file_info['path']
    file_name = file_info['name']
    file_type = file_info['type']
    
    all_text_chunks = []
    
    try:
        # First, add the file name as a searchable chunk
        file_name_chunk = {
            "file_name": file_name,
            "page_number": 0,  # Special page number for file name
            "text": f"Filename: {file_name}",
            "file_type": file_type,
            "chunk_type": "filename"
        }
        all_text_chunks.append(file_name_chunk)
        
        # Then extract content based on file type
        content_chunks = []
        if file_type == 'pdf':
            content_chunks = extract_pdf_text(file_path, file_name)
        elif file_type in ['docx', 'doc']:
            content_chunks = extract_docx_text(file_path, file_name)
        elif file_type in ['xlsx', 'xls']:
            content_chunks = extract_excel_text(file_path, file_name)
        elif file_type == 'csv':
            content_chunks = extract_csv_text(file_path, file_name)
        elif file_type in ['pptx', 'ppt']:
            content_chunks = extract_pptx_text(file_path, file_name)
        elif file_type in ['txt', 'rtf', 'markdown', 'html', 'xml', 'json', 'log', 'ini', 'cfg', 'conf']:
            content_chunks = extract_text_file(file_path, file_name)
        elif file_type == 'image':
            content_chunks = extract_image_keywords(file_path, file_name)
        
        # For images, the chunks are already fully processed with specific types.
        # For other documents, we assign the generic 'content' type.
        if file_type == 'image':
            all_text_chunks.extend(content_chunks)
        else:
            for chunk in content_chunks:
                chunk["chunk_type"] = "content"
                all_text_chunks.append(chunk)
            
        return all_text_chunks
        
    except Exception as e:
        print(f"Error extracting text from {file_name}: {e}")
        return all_text_chunks  # Return at least the filename chunk

def extract_pdf_text(pdf_path, file_name):
    """Extract text from PDF files"""
    reader = PdfReader(pdf_path)
    all_text = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        all_text.append({
            "file_name": file_name,
            "page_number": i,
            "text": text.strip(),
            "file_type": "pdf"
        })
    return all_text

def extract_docx_text(docx_path, file_name):
    """Extract text from Word documents"""
    doc = Document(docx_path)
    all_text = []
    
    # Extract text from paragraphs
    full_text = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            full_text.append(paragraph.text.strip())
    
    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text.strip():
                    row_text.append(cell.text.strip())
            if row_text:
                full_text.append(" | ".join(row_text))
    
    # Combine all text
    combined_text = "\n".join(full_text)
    if combined_text.strip():
        all_text.append({
            "file_name": file_name,
            "page_number": 1,  # Word docs don't have pages like PDFs
            "text": combined_text,
            "file_type": "docx"
        })
    
    return all_text

def extract_excel_text(excel_path, file_name):
    """Extract text from Excel files"""
    workbook = openpyxl.load_workbook(excel_path, data_only=True)
    all_text = []
    
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        sheet_text = []
        
        for row in sheet.iter_rows():
            row_text = []
            for cell in row:
                if cell.value is not None:
                    row_text.append(str(cell.value).strip())
            if row_text:
                sheet_text.append(" | ".join(row_text))
        
        if sheet_text:
            combined_text = "\n".join(sheet_text)
            all_text.append({
                "file_name": f"{file_name} - {sheet_name}",
                "page_number": 1,
                "text": combined_text,
                "file_type": "excel"
            })
    
    return all_text

def extract_csv_text(csv_path, file_name):
    """Extract text from CSV files"""
    all_text = []
    
    with open(csv_path, 'r', encoding='utf-8', errors='ignore') as file:
        csv_reader = csv.reader(file)
        rows = []
        
        for row in csv_reader:
            if any(cell.strip() for cell in row):  # Skip empty rows
                rows.append(" | ".join(cell.strip() for cell in row if cell.strip()))
        
        if rows:
            combined_text = "\n".join(rows)
            all_text.append({
                "file_name": file_name,
                "page_number": 1,
                "text": combined_text,
                "file_type": "csv"
            })
    
    return all_text

def extract_pptx_text(pptx_path, file_name):
    """Extract text from PowerPoint presentations"""
    prs = Presentation(pptx_path)
    all_text = []
    
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        # Extract text from tables
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        slide_text.append(" | ".join(row_text))
        
        if slide_text:
            combined_text = "\n".join(slide_text)
            all_text.append({
                "file_name": f"{file_name} - Slide {i}",
                "page_number": i,
                "text": combined_text,
                "file_type": "pptx"
            })
    
    return all_text

def extract_text_file(file_path, file_name):
    """Extract text from plain text files"""
    all_text = []
    
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    content = file.read().strip()
                    if content:
                        all_text.append({
                            "file_name": file_name,
                            "page_number": 1,
                            "text": content,
                            "file_type": "text"
                        })
                        break
            except UnicodeDecodeError:
                continue
                
    except Exception as e:
        print(f"Error reading text file {file_name}: {e}")
    
    return all_text

def extract_all_documents(folder_path):
    """Extract text from all supported documents in a folder"""
    documents = find_documents(folder_path)
    all_extracted_text = []
    
    print(f"Found {len(documents)} supported documents")
    
    for doc in documents:
        print(f"Processing: {doc['name']} ({doc['type']})")
        extracted_text = extract_text_from_file(doc)
        all_extracted_text.extend(extracted_text)
    
    return all_extracted_text

# Backward compatibility functions
def find_pdfs(folder_path):
    """Backward compatibility - find only PDF files"""
    documents = find_documents(folder_path)
    return [doc['path'] for doc in documents if doc['type'] == 'pdf']

def extract_image_keywords(image_path, file_name):
    """Extract keywords from images using AI classification"""
    all_text = []
    
    try:
        # Get the image classifier
        classifier = get_image_classifier()
        
        if classifier is None:
            # Fallback: create basic description from filename
            print(f"⚠️ Image classifier not available, using filename for {file_name}")
            basic_description = f"Image file: {os.path.splitext(file_name)[0]}"
            all_text.append({
                "file_name": file_name,
                "page_number": 1,
                "text": basic_description,
                "file_type": "image",
                "chunk_type": "image_description"
            })
            return all_text
        
        # Use AI to classify the image and generate keywords
        print(f"🔍 Analyzing image: {file_name}")
        predictions = classifier.predict_from_path(image_path, top_k=10)  # Get more predictions
        
        # Create comprehensive searchable text from predictions
        keywords = []
        high_confidence_items = []
        medium_confidence_items = []
        all_items = []
        
        for i, (class_name, confidence) in enumerate(predictions):
            if confidence > 0.05:  # Include more predictions for better coverage
                # Clean up class names (remove underscores, make more natural)
                clean_name = class_name.replace('_', ' ').replace('-', ' ')
                keywords.append(clean_name)
                all_items.append(clean_name)
                
                if confidence > 0.3:  # High confidence
                    high_confidence_items.append(clean_name)
                elif confidence > 0.1:  # Medium confidence
                    medium_confidence_items.append(clean_name)
        
        # Create multiple searchable text chunks for better matching
        
        # 1. Filename chunk
        all_text.append({
            "file_name": file_name,
            "page_number": 1,
            "text": f"Image filename: {file_name}",
            "file_type": "image",
            "chunk_type": "image_filename"
        })
        
        # 2. Natural description chunk - most important for search
        if high_confidence_items:
            natural_description = f"This image shows {', '.join(high_confidence_items[:3])}"
            all_text.append({
                "file_name": file_name,
                "page_number": 1,
                "text": natural_description,
                "file_type": "image",
                "chunk_type": "image_description"
            })
        
        # 3. Comprehensive content description
        if all_items:
            comprehensive_description = f"Image contains: {', '.join(all_items[:8])}"
            all_text.append({
                "file_name": file_name,
                "page_number": 1,
                "text": comprehensive_description,
                "file_type": "image",
                "chunk_type": "image_content"
            })
        
        # 4. Individual high-confidence items as separate chunks for precise matching
        for item in high_confidence_items[:5]:
            # Create natural language descriptions
            searchable_texts = [
                f"This is a {item}",
                f"Image of {item}",
                f"Contains {item}",
                f"Shows {item}",
                item  # Just the item name itself
            ]
            
            for searchable_text in searchable_texts:
                all_text.append({
                    "file_name": file_name,
                    "page_number": 1,
                    "text": searchable_text,
                    "file_type": "image",
                    "chunk_type": "image_searchable"
                })
        
        # 5. Medium confidence items for broader search coverage
        for item in medium_confidence_items[:3]:
            all_text.append({
                "file_name": file_name,
                "page_number": 1,
                "text": f"Possibly contains {item}",
                "file_type": "image",
                "chunk_type": "image_possible"
            })
        
        print(f"✅ Generated {len(all_text)} searchable chunks for image: {file_name}")
        
    except Exception as e:
        print(f"Error processing image {file_name}: {e}")
        # Fallback: create basic description from filename
        all_text.append({
            "file_name": file_name,
            "page_number": 1,
            "text": f"Image file: {os.path.splitext(file_name)[0]}",
            "file_type": "image",
            "chunk_type": "image_description"
        })
    
    return all_text

def extract_pdf_text_legacy(pdf_path):
    """Backward compatibility - extract text from PDF"""
    file_name = os.path.basename(pdf_path)
    return extract_pdf_text(pdf_path, file_name)